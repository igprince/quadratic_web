# app.py (fixed version)
import io
import base64
import logging
from flask import Flask, request, render_template, send_file, redirect, flash, url_for
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os




app = Flask(__name__)
app.secret_key = "supersecretkey"
logging.basicConfig(level=logging.INFO)

def calculate_and_explain(a, b, c, as_html=True):
    D = b**2 - 4*a*c
    sqrt_part = np.sqrt(abs(D))
    vertex_x = -b / (2*a)
    vertex_y = a * vertex_x**2 + b * vertex_x + c

    if D > 0:
        nature = "real and distinct"
        r1 = (-b + sqrt_part) / (2*a)
        r2 = (-b - sqrt_part) / (2*a)
        roots = f"x = (-{b} ± √{D}) / 2({a}) = ±{sqrt_part:.2f} / {2*a} = {r1:.2f}, {r2:.2f}"
    elif D == 0:
        nature = "real and equal"
        r1 = r2 = -b / (2*a)
        roots = f"x = -{b} / 2({a}) = {r1:.2f}"
    else:
        nature = "complex"
        real = -b / (2*a)
        imag = sqrt_part / (2*a)
        roots = f"x = {real:.2f} ± {imag:.2f}i"
        r1 = r2 = None

    vertex_calc = f"x = -{b}/(2×{a}) = {vertex_x:.2f}, y = {a}({vertex_x:.2f})² + {b}({vertex_x:.2f}) + {c} = {vertex_y:.2f}"

    explanation_txt = [
        f"Equation: y = {a}x² + {b}x + {c}",
        f"Discriminant (D) = {b}² - 4({a})({c}) = {D}",
        f"Since D {'>' if D > 0 else '=' if D == 0 else '<'} 0, roots are {nature}.",
        "x = (-b ± √D) / 2a",
        roots,
        f"Vertex: {vertex_calc}"
    ]

    img_stream = io.BytesIO()
    img_base64 = ""
    try:
        x_vals = np.linspace(vertex_x - 10, vertex_x + 10, 400)
        y_vals = a * x_vals**2 + b * x_vals + c
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.plot(x_vals, y_vals, label=f"y = {a}x² + {b}x + {c}", color='navy')
        ax.axhline(0, color='black', linewidth=0.8)
        ax.axvline(0, color='black', linewidth=0.8)
        ax.grid(True, linestyle='--', linewidth=0.5, alpha=0.6)
        ax.scatter(vertex_x, vertex_y, color='red', label='Vertex')
        if D >= 0:
            ax.scatter([r1, r2], [0, 0], color='green', label='Roots')
        ax.scatter(0, c, color='purple', label='Y-intercept')
        ax.set_title("Graph of Quadratic Function")
        ax.legend()
        fig.tight_layout()
        fig.savefig(img_stream, format='png', dpi=140)
        plt.close(fig)
        img_stream.seek(0)
        img_base64 = base64.b64encode(img_stream.getvalue()).decode('utf-8')
    except Exception as e:
        logging.error(f"Graph error: {e}")

    if as_html:
        explanation_html = "<br><br>".join(f"<pre>{line}</pre>" for line in explanation_txt)
        return explanation_html, img_base64, img_stream
    else:
        return explanation_txt, img_base64, img_stream

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'GET':
        return render_template("index.html", explanation="", img_base64="", a_val="", b_val="", c_val="")

    try:
        a = float(request.form['a'])
        b = float(request.form['b'])
        c = float(request.form['c'])
        if a == 0:
            flash("Coefficient 'a' must not be zero.", "error")
            return redirect(url_for('index'))
    except (ValueError, KeyError):
        flash("Please enter valid numerical values for a, b, and c.", "error")
        return redirect(url_for('index'))

    explanation, img_base64, _ = calculate_and_explain(a, b, c, as_html=True)

    return render_template("index.html", explanation=explanation, img_base64=img_base64, a_val=a, b_val=b, c_val=c)

@app.route('/download_ppt', methods=['POST'])
def download_ppt():
    a = float(request.form['a'])
    b = float(request.form['b'])
    c = float(request.form['c'])

    explanation_txt, _, img_stream = calculate_and_explain(a, b, c, as_html=False)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def rgb(r, g, b): return RGBColor(r, g, b)
    main_color = rgb(36, 41, 78)
    light_bg = rgb(240, 245, 255)
    dark_bg = rgb(36, 41, 78)
    highlight_color = rgb(255, 209, 102)
    white = rgb(255, 255, 255)
    theme_font = "Calibri"

    def add_slide(title, body_lines, bg_color=light_bg, title_color=main_color, body_color=main_color):
        slide = prs.slides.add_slide(blank)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = title_color
        p.font.name = theme_font

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5))
        tfb = body_box.text_frame
        for line in body_lines:
            p = tfb.add_paragraph()
            p.text = line
            p.font.size = Pt(22)
            p.font.name = theme_font
            p.font.color.rgb = body_color

    # Intro slide
    intro = prs.slides.add_slide(blank)
    bg = intro.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = dark_bg
    title_box = intro.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title = title_box.text_frame
    title.text = "📘 Quadratic Equation Visualizer"
    title.paragraphs[0].font.size = Pt(48)
    title.paragraphs[0].font.color.rgb = white
    title.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_box = intro.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    subtitle = subtitle_box.text_frame
    subtitle.text = "Created with Python, Math & Magic ✨"
    subtitle.paragraphs[0].font.size = Pt(24)
    subtitle.paragraphs[0].font.color.rgb = rgb(170, 220, 255)
    subtitle.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_box = intro.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    footer = footer_box.text_frame
    footer.text = "By IG_PRINCE"
    footer.paragraphs[0].font.size = Pt(20)
    footer.paragraphs[0].font.color.rgb = rgb(180, 180, 255)
    footer.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Content slides
    add_slide("🎯 What is a Quadratic Equation?", [
        "y = ax² + bx + c where a ≠ 0",
        "Forms a parabola when graphed",
        "Discriminant D = b² - 4ac tells root nature"
    ])
    add_slide("🧠 Steps & Solution", explanation_txt)

    # Graph slide
    slide = prs.slides.add_slide(blank)
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    tf = title.text_frame
    tf.text = "📈 Graph of the Equation"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.color.rgb = main_color
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), Inches(7.5), Inches(4.5))

    # Summary + Outro
    add_slide("💡 What Did We Learn?", [
        "• Quadratic equations create parabolas",
        "• The vertex is a key turning point",
        "• Discriminant tells the nature of roots",
        "• Python + Flask can visualize everything!"
    ])
    add_slide("🙏 Thank You!", [
        "Made with ❤️ using Python.",
        "Explore. Visualize. Learn."
    ], bg_color=dark_bg, title_color=highlight_color, body_color=white)

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return send_file(ppt_stream, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                     as_attachment=True, download_name="Quadratic_Eq_Best_Presentation.pptx")
@app.route('/download_png', methods=['POST'])
def download_png():
    try:
        a = float(request.form['a'])
        b = float(request.form['b'])
        c = float(request.form['c'])
        _, _, img_stream = calculate_and_explain(a, b, c, as_html=False)
        img_stream.seek(0)
        return send_file(img_stream, mimetype='image/png',
                         as_attachment=True, download_name='quadratic_graph.png')
    except Exception as e:
        logging.error(f"PNG download error: {e}")
        flash("Something went wrong while generating the PNG.", "error")
        return redirect(url_for('index'))

from flask import Response

@app.route('/download_pdf', methods=['POST'])

@app.route('/download_pdf', methods=['POST'])
def download_pdf():
    try:
        from fpdf import FPDF

        # 1) Read coefficients
        a = float(request.form['a'])
        b = float(request.form['b'])
        c = float(request.form['c'])

        # 2) Get explanation text and graph stream
        explanation_txt, _, img_stream = calculate_and_explain(a, b, c, as_html=False)
        img_data = img_stream.getvalue()

        # 3) Write image to a temp file (fpdf requires a filename for images)
        temp_img = 'temp_graph.png'
        with open(temp_img, 'wb') as f:
            f.write(img_data)

        # 4) Build PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in explanation_txt:
            pdf.multi_cell(0, 10, line)
        pdf.image(temp_img, x=10, w=180)

        # 5) Clean up the temp image
        os.remove(temp_img)

        # 6) Output PDF as bytes
        pdf_str = pdf.output(dest='S')               # get PDF as string
        pdf_bytes = pdf_str.encode('latin-1')         # encode to bytes

        # 7) Send it directly
        return Response(
            pdf_bytes,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': 'attachment; filename="quadratic_summary.pdf"'
            }
        )

    except Exception:
        logging.exception("PDF download error")
        flash("Oops, something went wrong generating the PDF.", "error")
        return redirect(url_for('index'))




if __name__ == '__main__':
    app.run(debug=True, host="0.0.0.0", port=5000)
    # Note: In a production environment, set debug=False and use a proper WSGI server like Gunicorn or uWSGI.   